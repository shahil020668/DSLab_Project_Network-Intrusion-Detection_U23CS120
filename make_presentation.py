from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT = Path(__file__).resolve().parent / "Network_Intrusion_Detection_Project_Presentation.pptx"

TITLE_COLOR = RGBColor(11, 61, 145)
ACCENT = RGBColor(0, 102, 204)
TEXT = RGBColor(40, 40, 40)
MUTED = RGBColor(90, 90, 90)
BG = RGBColor(248, 250, 252)


slides = [
    {
        "title": "Network Intrusion Detection System",
        "subtitle": "Data Science Lab Project\nNSL-KDD Dataset\n\nTeam Members: [Edit names]\nRoll No: [Edit roll numbers]\nGuide: [Edit guide name]",
        "bullets": []
    },
    {
        "title": "Problem Statement and Motivation",
        "subtitle": "Why this problem matters",
        "bullets": [
            "Goal: detect malicious network traffic from normal traffic.",
            "Cyber attacks can be hidden in high-volume traffic streams.",
            "Manual monitoring is slow and error-prone.",
            "Machine learning helps detect attacks earlier and more reliably.",
            "[Edit this slide to match your exact project motivation.]"
        ]
    },
    {
        "title": "Dataset Description",
        "subtitle": "NSL-KDD dataset overview",
        "bullets": [
            "Dataset source: NSL-KDD.",
            "Input features: 41 network-traffic features.",
            "Target label: Normal vs Attack.",
            "Used for classification and anomaly detection.",
            "[Insert dataset statistics, sample records, and class distribution here.]"
        ]
    },
    {
        "title": "Data Preprocessing",
        "subtitle": "Cleaning and transformation steps",
        "bullets": [
            "Encoded categorical features such as protocol_type, service, and flag.",
            "Scaled features using MinMaxScaler.",
            "Handled imbalance using SMOTE.",
            "Split data into training and testing sets.",
            "[Add a pipeline diagram or preprocessing flowchart.]"
        ]
    },
    {
        "title": "Exploratory Data Analysis",
        "subtitle": "Visual patterns in the data",
        "bullets": [
            "Show class distribution (normal vs attack).",
            "Show feature importance or correlation plots.",
            "Highlight key trends observed in the dataset.",
            "[Insert charts from your notebooks/output folder.]"
        ]
    },
    {
        "title": "Methodology",
        "subtitle": "Models and workflow",
        "bullets": [
            "Classification models: Random Forest, XGBoost, SVM, Logistic Regression.",
            "Anomaly detection: IsolationForest.",
            "Feature selection: Extra Trees, PCA, and RFE.",
            "Hybrid decision logic combines model outputs.",
            "[Add a workflow diagram here.]"
        ]
    },
    {
        "title": "Results and Evaluation",
        "subtitle": "Model performance summary",
        "bullets": [
            "Report Accuracy, Precision, Recall, F1-score, and ROC-AUC.",
            "Highlight the best-performing model.",
            "Emphasize Recall because missing attacks is costly.",
            "[Paste your final comparison table or metric chart here.]"
        ]
    },
    {
        "title": "Comparison with Existing Work",
        "subtitle": "Reference paper comparison",
        "bullets": [
            "Compare your method with the reference papers.",
            "Show what you improved: hybrid ensemble, anomaly detection, explainability.",
            "Mention whether results improved on Recall or F1-score.",
            "[Add a side-by-side comparison table.]"
        ]
    },
    {
        "title": "Challenges, Limitations, Future Work",
        "subtitle": "What we learned and what can be improved",
        "bullets": [
            "Challenges: imbalance, dimensionality, and model tuning.",
            "Limitations: dataset is simulated/benchmark traffic, not live traffic.",
            "Future work: real-time deployment, API integration, multiclass attacks.",
            "[Edit to match your actual implementation experience.]"
        ]
    },
    {
        "title": "Conclusion / Demo",
        "subtitle": "Final takeaway",
        "bullets": [
            "The system detects malicious traffic using ML, DL, and anomaly detection.",
            "Feature selection and preprocessing improve reliability.",
            "Streamlit dashboard can be used for live demo.",
            "Thank you.",
            "[Replace this with your final summary and demo screenshots.]"
        ]
    },
]


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.0), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.color.rgb = MUTED


def add_footer(slide, number):
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(2.0), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"Slide {number}"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Remove default slides if any
while len(prs.slides) > 0:
    r_id = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[0]

for i, item in enumerate(slides, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, item["title"], item["subtitle"])

    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(11.7), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    if i == 1:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = item["subtitle"]
        r.font.size = Pt(20)
        r.font.color.rgb = TEXT
        r.font.bold = True
    else:
        for idx, bullet in enumerate(item["bullets"]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT
            p.bullet = True

    # Add editable placeholder box for visuals
    placeholder = slide.shapes.add_shape(1, Inches(8.7), Inches(4.8), Inches(3.7), Inches(1.3))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(255, 255, 255)
    placeholder.line.color.rgb = ACCENT
    tx = placeholder.text_frame
    tx.clear()
    p = tx.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Insert chart / screenshot here"
    r.font.size = Pt(12)
    r.font.color.rgb = ACCENT
    r.font.bold = True

    add_footer(slide, i)

prs.save(str(OUTPUT))
print(f"Saved presentation to: {OUTPUT}")
